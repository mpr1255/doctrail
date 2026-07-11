"""Extraction helpers for Microsoft PowerPoint presentations."""

from __future__ import annotations

import logging
import shutil
import subprocess
import tempfile
from pathlib import Path

logger = logging.getLogger(__name__)


def _resolve_soffice() -> str | None:
    """Return the LibreOffice CLI path if available."""
    return shutil.which("soffice") or shutil.which("libreoffice")


def _convert_with_soffice(file_path: str, target_ext: str) -> str:
    """Convert an Office document using LibreOffice and return the converted path."""
    soffice = _resolve_soffice()
    if not soffice:
        raise RuntimeError(
            f"LibreOffice is not installed. Install `soffice` to convert {Path(file_path).suffix} files."
        )

    with tempfile.TemporaryDirectory(prefix="doctrail-office-") as tmpdir:
        cmd = [
            soffice,
            "--headless",
            "--nologo",
            "--nodefault",
            "--nolockcheck",
            "--convert-to",
            target_ext,
            "--outdir",
            tmpdir,
            file_path,
        ]
        logger.debug("Running soffice: %s", " ".join(cmd))
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=180)

        converted_path = Path(tmpdir) / f"{Path(file_path).stem}.{target_ext}"
        if result.returncode != 0 or not converted_path.exists():
            stderr = result.stderr.strip() or result.stdout.strip()
            raise RuntimeError(
                f"LibreOffice failed to convert {Path(file_path).name} to {target_ext}: {stderr}"
            )

        persistent_tmpdir = tempfile.mkdtemp(prefix="doctrail-office-converted-")
        persistent_path = Path(persistent_tmpdir) / converted_path.name
        shutil.copy2(converted_path, persistent_path)
        return str(persistent_path)


def _collect_shape_lines(shape) -> list[str]:
    """Extract readable text from a PowerPoint shape."""
    lines: list[str] = []

    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            text = paragraph.text.strip()
            if text:
                lines.append(text)

    if hasattr(shape, "has_table") and shape.has_table:
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            while cells and not cells[-1]:
                cells.pop()
            if any(cells):
                lines.append(" | ".join(cells))

    if hasattr(shape, "shapes"):
        for child_shape in shape.shapes:
            lines.extend(_collect_shape_lines(child_shape))

    return lines


def extract_text_from_pptx(file_path: str) -> tuple[str, str]:
    """
    Extract text from a .pptx presentation.

    Returns:
        (content, extraction_method)
    """
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx is not installed. Run `uv sync` to ingest .pptx files.") from exc

    presentation = Presentation(file_path)
    slide_sections: list[str] = []

    for index, slide in enumerate(presentation.slides, start=1):
        seen: set[str] = set()
        slide_lines: list[str] = []

        for shape in slide.shapes:
            for line in _collect_shape_lines(shape):
                if line not in seen:
                    seen.add(line)
                    slide_lines.append(line)

        if slide_lines:
            slide_sections.append(f"Slide {index}\n" + "\n".join(slide_lines))

    return "\n\n".join(slide_sections).strip(), "python_pptx"


def _extract_text_from_ppt_with_strings(file_path: str) -> str:
    strings = shutil.which("strings")
    if not strings:
        return ""

    result = subprocess.run(
        [strings, "-n", "4", file_path],
        capture_output=True,
        text=True,
        timeout=60,
    )
    if result.returncode != 0:
        return ""

    ignored = {
        "PowerPoint Document",
        "PowerPoint Presentation",
        "Microsoft PowerPoint",
        "Current User",
        "Pictures",
        "Document",
    }
    lines: list[str] = []
    seen: set[str] = set()
    for raw_line in result.stdout.splitlines():
        line = " ".join(raw_line.split())
        if not line or line in ignored:
            continue
        if len(line) < 4 or not any(ch.isalpha() for ch in line):
            continue
        if line not in seen:
            seen.add(line)
            lines.append(line)

    return "\n".join(lines).strip()


def extract_text_from_ppt(file_path: str) -> tuple[str, str]:
    """
    Extract text from a legacy .ppt presentation.

    Returns:
        (content, extraction_method)
    """
    content = _extract_text_from_ppt_with_strings(file_path)
    if content:
        return content, "strings"

    converted_path = None
    try:
        converted_path = _convert_with_soffice(file_path, "pptx")
        content, _ = extract_text_from_pptx(converted_path)
        return content, "soffice_to_pptx_python_pptx"
    finally:
        if converted_path:
            converted = Path(converted_path)
            try:
                converted.unlink(missing_ok=True)
                converted.parent.rmdir()
            except OSError:
                pass
