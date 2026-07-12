import hashlib
import io
import json
import shutil
import sqlite3
import subprocess
import zipfile
from pathlib import Path
from types import SimpleNamespace

import pytest
import respx
from httpx import Response

from doctrail.ingest.document_processor import (
    SUPPORTED_EXTENSION_FAMILIES,
    SkippedFileException,
    _load_mac_ocr_module,
    _ocr_with_mac_ocr_service,
    format_supported_extensions_for_help,
    process_document,
)
from doctrail.ingest import embedded_media
from doctrail.ingest.core import process_ingest
from doctrail.ingest.text_processing import clean_extracted_text
from doctrail.extractors import spreadsheet_extractor

ASSET_DIR = Path(__file__).parent / "assets" / "files"


def _sha1_for(path):
    return hashlib.sha1(path.read_bytes()).hexdigest()


def _async_return(value):
    async def _inner(*args, **kwargs):
        return value
    return _inner


def _missing_tools(*tools):
    return [tool for tool in tools if shutil.which(tool) is None]


@pytest.mark.parametrize(
    ("family", "filename", "expected_methods", "minimum_chars", "expected_phrase", "required_tools"),
    [
        ("text", "federalist_fixture.txt", {"direct_text_read"}, 150, "Federalist fixture", ()),
        ("text", "federalist_fixture.md", {"direct_text_read"}, 150, "Federalist fixture", ()),
        ("delimited", "federalist_fixture.csv", {"csv"}, 80, "Federalist No. 1", ()),
        ("delimited", "federalist_fixture.tsv", {"csv"}, 80, "Federalist No. 1", ()),
        ("pdf", "federalist_fixture.pdf", {"pymupdf", "pdftotext", "mutool"}, 150, "Federalist fixture", ()),
        ("epub", "federalist_fixture.epub", {"epub_direct"}, 150, "Federalist fixture", ()),
        ("mobi", "federalist_fixture.mobi", {"ebook_convert"}, 150, "Federalist fixture", ("ebook-convert",)),
        ("word", "federalist_fixture.doc", {"antiword"}, 150, "Federalist fixture", ("antiword",)),
        ("rtf", "federalist_fixture.rtf", {"textutil"}, 150, "Federalist fixture", ("textutil",)),
        ("word_openxml", "federalist_fixture.docx", {"python_docx"}, 150, "Federalist fixture", ()),
        ("excel_openxml", "federalist_fixture.xlsx", {"openpyxl"}, 80, "Federalist No. 1", ()),
        ("excel", "federalist_fixture.xls", {"xlrd"}, 80, "Federalist No. 1", ()),
        ("powerpoint_openxml", "federalist_fixture.pptx", {"python_pptx"}, 80, "Federalist fixture", ()),
        ("powerpoint", "federalist_fixture.ppt", {"soffice_to_pptx_python_pptx", "strings"}, 80, "Federalist fixture", ("strings",)),
        ("djvu", "federalist_fixture.djvu", {"djvutxt"}, 80, "Federalist fixture", ("djvutxt",)),
        ("mhtml", "federalist_fixture.mhtml", {"mhtml_header_parsing", "beautifulsoup"}, 150, "Federalist fixture", ()),
        ("html", "federalist_fixture.html", {"beautifulsoup"}, 150, "Federalist fixture", ()),
        ("image_ocr", "federalist_fixture.png", {"textra_ocr", "tesseract_ocr"}, 40, "Federalist fixture", ("tesseract",)),
    ],
)
@pytest.mark.asyncio
async def test_process_document_handles_committed_real_file_fixtures(
    family,
    filename,
    expected_methods,
    minimum_chars,
    expected_phrase,
    required_tools,
):
    missing = _missing_tools(*required_tools)
    if missing:
        pytest.skip(f"missing required extraction tool(s): {', '.join(missing)}")

    path = ASSET_DIR / filename
    assert path.exists(), f"missing committed fixture for {family}: {filename}"
    assert path.suffix.lower() in SUPPORTED_EXTENSION_FAMILIES[family]

    sha1 = _sha1_for(path)
    result_sha1, content, metadata = await process_document(str(path), sha1)

    assert result_sha1 == sha1
    assert isinstance(content, str)
    assert len(content) >= minimum_chars
    assert expected_phrase in content
    assert metadata["extraction_method"] in expected_methods
    assert metadata["original_file_path"] == str(path)
    assert metadata["original_file_type"] == path.suffix.lower().lstrip(".")
    assert metadata["resourceName"] == filename
    if "title" in metadata:
        assert isinstance(metadata["title"], str)


def test_supported_extension_help_uses_dispatch_source_of_truth():
    help_text = format_supported_extensions_for_help()

    for extensions in SUPPORTED_EXTENSION_FAMILIES.values():
        grouped = "/".join(extension.lstrip(".") for extension in extensions)
        assert grouped in help_text


def test_load_mac_ocr_module_from_env_path(monkeypatch, tmp_path):
    client_path = tmp_path / "local_ocr_client.py"
    client_path.write_text(
        "async def ocr_async(file_path, node=None):\n"
        "    return 'local ocr text'\n",
        encoding="utf-8",
    )

    monkeypatch.setenv("DOCTRAIL_MAC_OCR_CLIENT_PATH", str(client_path))
    module = _load_mac_ocr_module()

    assert hasattr(module, "ocr_async")


@pytest.mark.asyncio
@respx.mock
async def test_mac_ocr_service_uses_configured_funnel_reservation(monkeypatch, tmp_path):
    first = "https://ocr-one.example/funnel"
    second = "https://ocr-two.example/funnel"
    monkeypatch.setenv("MAC_OCR__SERVICE_ENDPOINTS", f"{first},{second}")
    image_path = tmp_path / "page.png"
    image_path.write_bytes(b"image bytes")

    respx.post(f"{first}/reserve").mock(return_value=Response(503))
    respx.post(f"{second}/reserve").mock(
        return_value=Response(201, json={"reservation_id": "reserved"})
    )
    upload = respx.post(f"{second}/ocr", params={"reservation_id": "reserved"}).mock(
        return_value=Response(200, json={"text": "==== Page 1 ====\nFarm OCR text"})
    )

    text = await _ocr_with_mac_ocr_service(str(image_path))

    assert "Farm OCR text" in text
    assert upload.called


@pytest.mark.asyncio
@respx.mock
@pytest.mark.parametrize("status_code", [422, 500])
async def test_mac_ocr_service_reports_file_rejection_without_capacity_loop(
    monkeypatch, tmp_path, status_code
):
    endpoint = "https://ocr-one.example/funnel"
    monkeypatch.setenv("MAC_OCR__SERVICE_ENDPOINTS", endpoint)
    image_path = tmp_path / "broken.pdf"
    image_path.write_bytes(b"%PDF-broken")

    respx.post(f"{endpoint}/reserve").mock(
        return_value=Response(201, json={"reservation_id": "reserved"})
    )
    respx.post(f"{endpoint}/ocr", params={"reservation_id": "reserved"}).mock(
        return_value=Response(status_code, text="unsupported or unrenderable file")
    )

    with pytest.raises(
        RuntimeError,
        match=rf"HTTP {status_code}.*unsupported or unrenderable file",
    ):
        await _ocr_with_mac_ocr_service(str(image_path))


@pytest.mark.asyncio
async def test_process_document_handles_xlsx_real(tmp_path):
    from openpyxl import Workbook

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Summary"
    sheet.append(["Name", "Count"])
    sheet.append(["Alice", 3])
    sheet.append(["Bob", 4])

    xlsx_path = tmp_path / "sample.xlsx"
    workbook.save(xlsx_path)

    sha1 = _sha1_for(xlsx_path)
    result_sha1, content, metadata = await process_document(str(xlsx_path), sha1)

    assert result_sha1 == sha1
    assert "Sheet: Summary" in content
    assert "Alice | 3" in content
    assert metadata["original_file_path"] == str(xlsx_path)
    assert metadata["original_file_type"] == "xlsx"
    assert metadata["resourceName"] == "sample.xlsx"
    assert metadata["Content-Type"] == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    assert metadata["extraction_method"] == "openpyxl"


@pytest.mark.asyncio
async def test_process_document_handles_docx_real(tmp_path):
    from docx import Document

    document = Document()
    document.core_properties.title = "Federalist fixture"
    document.add_heading("Federalist fixture", level=1)
    document.add_paragraph("This document has enough real text to catch tuple contract regressions.")

    docx_path = tmp_path / "sample.docx"
    document.save(docx_path)

    sha1 = _sha1_for(docx_path)
    result_sha1, content, metadata = await process_document(str(docx_path), sha1)

    assert result_sha1 == sha1
    assert "Federalist fixture" in content
    assert "tuple contract regressions" in content
    assert metadata["original_file_path"] == str(docx_path)
    assert metadata["original_file_type"] == "docx"
    assert metadata["resourceName"] == "sample.docx"
    assert metadata["Content-Type"] == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    assert metadata["extraction_method"] == "python_docx"
    assert metadata["title"] == "Federalist fixture"


def test_clean_extracted_text_rejects_non_string_input():
    with pytest.raises(TypeError):
        clean_extracted_text(("content", "title"))


def test_embedded_image_rows_preserve_each_parent_relationship(monkeypatch, tmp_path):
    first = tmp_path / "first.docx"
    second = tmp_path / "second.docx"
    first.write_bytes(b"first parent")
    second.write_bytes(b"second parent")
    shared_image = b"same screenshot bytes"

    monkeypatch.setattr(
        embedded_media,
        "office_images",
        lambda path, soffice="soffice": [("word/media/image1.png", shared_image)],
    )

    first_row = embedded_media._extract_rows_for_file(
        first, soffice="soffice", ocr_fn=lambda path: "", ocr_engine="none"
    )[0]
    second_row = embedded_media._extract_rows_for_file(
        second, soffice="soffice", ocr_fn=lambda path: "", ocr_engine="none"
    )[0]

    assert first_row["sha1"] != second_row["sha1"]
    assert first_row["extra_fields"]["image_sha1"] == second_row["extra_fields"]["image_sha1"]
    assert first_row["extra_fields"]["parent_sha1"] == _sha1_for(first)
    assert second_row["extra_fields"]["parent_sha1"] == _sha1_for(second)


def test_ooxml_embedded_media_rejects_zip_bomb_ratio():
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("word/media/image1.png", b"0" * (2 * 1024 * 1024))

    with pytest.raises(ValueError, match="compression ratio"):
        embedded_media._zip_media(buffer.getvalue())


def test_legacy_media_conversion_kills_process_group_on_timeout(monkeypatch, tmp_path):
    source = tmp_path / "legacy.doc"
    source.write_bytes(b"legacy")
    calls = {"communicate": 0}
    kill_calls = []

    class FakeProcess:
        pid = 12345
        returncode = None

        def communicate(self, timeout=None):
            calls["communicate"] += 1
            if calls["communicate"] == 1:
                raise subprocess.TimeoutExpired("soffice", timeout)
            return b"", b""

    def fake_popen(command, **kwargs):
        assert kwargs["start_new_session"] is True
        return FakeProcess()

    monkeypatch.setattr(embedded_media.subprocess, "Popen", fake_popen)
    monkeypatch.setattr(
        embedded_media.os,
        "killpg",
        lambda pid, sig: kill_calls.append((pid, sig)),
    )

    result = embedded_media._libreoffice_to_ooxml(
        source,
        "docx",
        tmp_path / "work",
        timeout=1,
    )

    assert result is None
    assert kill_calls == [(12345, embedded_media.signal.SIGKILL)]
    assert calls["communicate"] == 2


@pytest.mark.asyncio
async def test_native_ingest_extracts_stores_and_ocrs_embedded_docx_image(
    monkeypatch, tmp_path
):
    from docx import Document

    monkeypatch.delenv("DOCTRAIL_DISABLE_NATIVE", raising=False)
    source = tmp_path / "source"
    source.mkdir()
    image_path = ASSET_DIR / "federalist_fixture.png"
    image_bytes = image_path.read_bytes()
    docx_path = source / "with-image.docx"
    document = Document()
    document.add_paragraph("Parent document text for native extraction.")
    document.add_picture(str(image_path))
    document.save(docx_path)

    ocr_calls = []

    def fake_ocr(path):
        ocr_calls.append(path)
        return "Embedded image OCR fixture text"

    monkeypatch.setattr(
        embedded_media,
        "_ocr_callable",
        lambda engine: (fake_ocr, "test-ocr"),
    )
    db_path = tmp_path / "embedded.sqlite"
    result = await process_ingest(
        db_path=str(db_path),
        input_dir=str(source),
        table="documents",
        extractor="rust",
        ocr_engine="auto",
        workers=1,
        yes=True,
    )

    assert result["successful"] == 1
    assert result["failed"] == 0
    assert result["embedded_media"]["inserted"] == 1
    with sqlite3.connect(db_path) as connection:
        parent = connection.execute(
            "SELECT sha1, added_at FROM documents WHERE filename = 'with-image.docx'"
        ).fetchone()
        child = connection.execute(
            """SELECT parent_sha1, raw_content, image_bytes, metadata, added_at
               FROM documents WHERE source_format = 'embedded-image'"""
        ).fetchone()

    assert parent is not None
    assert child is not None
    assert child[0] == parent[0]
    assert child[1] == "Embedded image OCR fixture text"
    assert child[2] == image_bytes
    assert json.loads(child[3])["ocr_status"] == "success"
    assert parent[1]
    assert child[4]
    assert len(ocr_calls) == 1

    repeated = await process_ingest(
        db_path=str(db_path),
        input_dir=str(source),
        table="documents",
        extractor="rust",
        ocr_engine="auto",
        workers=1,
        yes=True,
    )
    assert repeated["successful"] == 0
    assert repeated["embedded_media"]["inserted"] == 0
    assert len(ocr_calls) == 1


@pytest.mark.asyncio
@pytest.mark.parametrize(
    ("extension", "extractor_name", "title", "expected_type"),
    [
        ("epub", "extract_text_from_epub", "EPUB title", "epub"),
        ("mobi", "extract_text_from_mobi", "MOBI title", "mobi"),
        ("djvu", "extract_text_from_djvu", "DJVU title", "djvu"),
    ],
)
async def test_process_document_unpacks_tuple_extractors(monkeypatch, tmp_path, extension, extractor_name, title, expected_type):
    file_path = tmp_path / f"sample.{extension}"
    file_path.write_bytes(b"fake document content")
    sha1 = _sha1_for(file_path)

    monkeypatch.setattr(
        f"doctrail.ingest.document_processor.{extractor_name}",
        lambda path: (f"{title}\nExtracted body text", title),
    )

    result_sha1, content, metadata = await process_document(str(file_path), sha1)

    assert result_sha1 == sha1
    assert content == f"{title}\nExtracted body text"
    assert metadata["original_file_type"] == expected_type
    assert metadata["title"] == title


@pytest.mark.asyncio
async def test_process_document_handles_xls(monkeypatch, tmp_path):
    xls_path = tmp_path / "sample.xls"
    xls_path.write_bytes(b"fake xls content")
    sha1 = _sha1_for(xls_path)

    monkeypatch.setattr(
        "doctrail.ingest.document_processor.extract_structured_from_xls",
        lambda path: (
            "Sheet: Legacy\nName | Count\nAlice | 5",
            "xlrd",
            [{
                "sheet_name": "Legacy",
                "rows": [["Name", "Count"], ["Alice", "5"]],
                "csv": "Name,Count\r\nAlice,5\r\n",
                "text": "Sheet: Legacy\nName | Count\nAlice | 5",
            }],
        ),
    )

    result_sha1, content, metadata = await process_document(str(xls_path), sha1)

    assert result_sha1 == sha1
    assert "Legacy" in content
    assert metadata["original_file_path"] == str(xls_path)
    assert metadata["original_file_type"] == "xls"
    assert metadata["resourceName"] == "sample.xls"
    assert metadata["Content-Type"] == "application/vnd.ms-excel"
    assert metadata["extraction_method"] == "xlrd"


@pytest.mark.asyncio
async def test_process_document_handles_html_backed_xls(tmp_path):
    xls_path = tmp_path / "sample.xls"
    xls_path.write_text(
        """<!DOCTYPE html>
<html>
<head><title>Workbook</title></head>
<body>
<table>
  <tr><td>Name</td><td>Count</td></tr>
  <tr><td>Alice</td><td>5</td></tr>
</table>
</body>
</html>
""",
        encoding="utf-8",
    )
    sha1 = _sha1_for(xls_path)

    result_sha1, content, metadata = await process_document(str(xls_path), sha1)

    assert result_sha1 == sha1
    assert "Sheet: Workbook" in content
    assert "Alice | 5" in content
    assert metadata["extraction_method"] == "html_spreadsheet"


@pytest.mark.asyncio
async def test_process_document_returns_plain_string_html_title(tmp_path):
    html_path = tmp_path / "sample.html"
    html_path.write_text(
        """<!DOCTYPE html>
<html>
<head><title>  Example   page  </title></head>
<body><h1>Visible heading</h1><p>Hello world</p></body>
</html>
""",
        encoding="utf-8",
    )
    sha1 = _sha1_for(html_path)

    _, content, metadata = await process_document(str(html_path), sha1)

    assert "Hello world" in content
    assert metadata["title"] == "Example page"
    assert isinstance(metadata["title"], str)


@pytest.mark.asyncio
async def test_process_document_sniffs_mhtml_saved_with_html_extension(tmp_path):
    html_path = tmp_path / "saved-page.html"
    shutil.copyfile(ASSET_DIR / "federalist_fixture.mhtml", html_path)
    sha1 = _sha1_for(html_path)

    result_sha1, content, metadata = await process_document(str(html_path), sha1)

    assert result_sha1 == sha1
    assert "Federalist fixture" in content
    assert metadata["file_type"] == "mhtml"


@pytest.mark.asyncio
async def test_process_document_handles_csv_real(tmp_path):
    csv_path = tmp_path / "sample.csv"
    csv_path.write_text(
        "Name,Count\nAlice,5\nBob,7\n",
        encoding="utf-8",
    )
    sha1 = _sha1_for(csv_path)

    result_sha1, content, metadata = await process_document(str(csv_path), sha1)

    assert result_sha1 == sha1
    assert "Sheet: sample" in content
    assert "Alice | 5" in content
    assert metadata["original_file_path"] == str(csv_path)
    assert metadata["original_file_type"] == "csv"
    assert metadata["resourceName"] == "sample.csv"
    assert metadata["Content-Type"] == "text/csv"
    assert metadata["extraction_method"] == "csv"


def test_soffice_retry_on_signal_terminated_process(monkeypatch, tmp_path):
    source_path = tmp_path / "sample.xls"
    source_path.write_bytes(b"fake xls content")

    monkeypatch.setattr(spreadsheet_extractor, "_resolve_soffice", lambda: "/usr/bin/soffice")

    class _Slot:
        def close(self):
            return None

    monkeypatch.setattr(spreadsheet_extractor, "_acquire_soffice_slot", lambda max_procs: _Slot())

    calls = {"count": 0}

    def _fake_run(cmd, capture_output, text, timeout):
        calls["count"] += 1
        outdir = Path(cmd[cmd.index("--outdir") + 1])
        converted_path = outdir / f"{source_path.stem}.xlsx"
        if calls["count"] == 2:
            converted_path.write_bytes(b"converted workbook")
            return SimpleNamespace(returncode=0, stderr="", stdout="")
        return SimpleNamespace(returncode=-2, stderr="", stdout="")

    monkeypatch.setattr(spreadsheet_extractor.subprocess, "run", _fake_run)
    monkeypatch.setenv("DOCTRAIL_SOFFICE_RETRIES", "1")

    converted_path = spreadsheet_extractor._convert_with_soffice(str(source_path), "xlsx")
    assert Path(converted_path).exists()
    assert calls["count"] == 2


@pytest.mark.asyncio
async def test_process_document_xls_missing_backend(monkeypatch, tmp_path):
    xls_path = tmp_path / "sample.xls"
    xls_path.write_bytes(b"fake xls content")
    sha1 = _sha1_for(xls_path)

    def _raise(_):
        raise RuntimeError("LibreOffice is not installed. Install `soffice` to convert .xls files.")

    monkeypatch.setattr("doctrail.ingest.document_processor.extract_structured_from_xls", _raise)

    with pytest.raises(ValueError) as excinfo:
        await process_document(str(xls_path), sha1)

    assert "soffice" in str(excinfo.value)


@pytest.mark.asyncio
async def test_process_document_skips_quota_html_saved_as_pdf(tmp_path):
    pdf_path = tmp_path / "quota.pdf"
    pdf_path.write_text(
        "<script>alert('对不起，今日的条目下载量已经达到极限，建议您明天下载！');</script>",
        encoding="utf-8",
    )
    sha1 = _sha1_for(pdf_path)

    with pytest.raises(SkippedFileException) as excinfo:
        await process_document(str(pdf_path), sha1)

    assert "Download-limit HTML page" in str(excinfo.value)


@pytest.mark.asyncio
async def test_process_document_skips_quota_html_saved_as_xls(tmp_path):
    xls_path = tmp_path / "quota.xls"
    xls_path.write_text(
        "<script>alert('对不起，今日的条目下载量已经达到极限，建议您明天下载！');</script>",
        encoding="utf-8",
    )
    sha1 = _sha1_for(xls_path)

    with pytest.raises(SkippedFileException) as excinfo:
        await process_document(str(xls_path), sha1)

    assert "Download-limit HTML page" in str(excinfo.value)


@pytest.mark.asyncio
async def test_process_document_respects_mac_ocr_engine(monkeypatch, tmp_path):
    pdf_path = tmp_path / "sample.pdf"
    pdf_path.write_bytes(b"%PDF-1.4\nfake pdf for engine selection\n")
    sha1 = _sha1_for(pdf_path)

    monkeypatch.setattr("doctrail.ingest.document_processor._try_ocr_with_mac_ocr", _async_return("==== Page 1 ====\nOCR text"))
    monkeypatch.setattr("doctrail.ingest.document_processor.extract_text_with_pymupdf", lambda path: "embedded text")

    result_sha1, content, metadata = await process_document(
        str(pdf_path),
        sha1,
        pdf_engine="mac-ocr",
        ocr_engine="mac-ocr",
    )

    assert result_sha1 == sha1
    assert "OCR text" in content
    assert metadata["extraction_method"] == "mac_ocr"
    assert metadata["ocr_engine"] == "mac-ocr"


@pytest.mark.asyncio
async def test_process_document_handles_pptx_real(tmp_path):
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Quarterly update"
    slide.placeholders[1].text = "Revenue up\nNew hires"

    pptx_path = tmp_path / "sample.pptx"
    presentation.save(pptx_path)

    sha1 = _sha1_for(pptx_path)
    result_sha1, content, metadata = await process_document(str(pptx_path), sha1)

    assert result_sha1 == sha1
    assert "Slide 1" in content
    assert "Quarterly update" in content
    assert "Revenue up" in content
    assert metadata["original_file_path"] == str(pptx_path)
    assert metadata["original_file_type"] == "pptx"
    assert metadata["resourceName"] == "sample.pptx"
    assert metadata["Content-Type"] == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    assert metadata["extraction_method"] == "python_pptx"


@pytest.mark.asyncio
async def test_process_document_handles_ppt(monkeypatch, tmp_path):
    ppt_path = tmp_path / "sample.ppt"
    ppt_path.write_bytes(b"fake ppt content")
    sha1 = _sha1_for(ppt_path)

    monkeypatch.setattr(
        "doctrail.ingest.document_processor.extract_text_from_ppt",
        lambda path: ("Slide 1\nLegacy deck", "soffice_to_pptx_python_pptx"),
    )

    result_sha1, content, metadata = await process_document(str(ppt_path), sha1)

    assert result_sha1 == sha1
    assert "Legacy deck" in content
    assert metadata["original_file_path"] == str(ppt_path)
    assert metadata["original_file_type"] == "ppt"
    assert metadata["resourceName"] == "sample.ppt"
    assert metadata["Content-Type"] == "application/vnd.ms-powerpoint"
    assert metadata["extraction_method"] == "soffice_to_pptx_python_pptx"


@pytest.mark.asyncio
async def test_process_document_ppt_missing_soffice(monkeypatch, tmp_path):
    ppt_path = tmp_path / "sample.ppt"
    ppt_path.write_bytes(b"fake ppt content")
    sha1 = _sha1_for(ppt_path)

    def _raise(_):
        raise RuntimeError("LibreOffice is not installed. Install `soffice` to convert .ppt files.")

    monkeypatch.setattr("doctrail.ingest.document_processor.extract_text_from_ppt", _raise)

    with pytest.raises(ValueError) as excinfo:
        await process_document(str(ppt_path), sha1)

    assert "soffice" in str(excinfo.value)
